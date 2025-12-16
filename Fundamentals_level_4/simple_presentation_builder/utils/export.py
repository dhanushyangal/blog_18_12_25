"""
PPTX export utility using python-pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


def create_pptx_from_screenshots(screenshots: list, output_file: str = "exports/presentation.pptx", presentation_title: str = "Presentation", progress_callback=None) -> str:
    """
    Create a PowerPoint file from screenshot images

    Args:
        screenshots: List of screenshot file paths
        output_file: Output PPTX file path
        presentation_title: Title of the presentation
        progress_callback: Optional callback function(event_type, data)

    Returns:
        Path to the created PPTX file
    """

    print(f"\n📊 Creating PPTX presentation...")

    # Create output directory
    output_path = Path(output_file)
    output_path.parent.mkdir(exist_ok=True, parents=True)

    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9 (standard widescreen)
    # Width: 10 inches, Height: 5.625 inches (16:9 ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Use blank slide layout (index 6)
    blank_slide_layout = prs.slide_layouts[6]

    for i, screenshot_path in enumerate(screenshots, 1):
        try:
            screenshot = Path(screenshot_path)

            if not screenshot.exists():
                print(f"   ⚠️  Screenshot not found: {screenshot_path}")
                continue

            # Add blank slide
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add screenshot to fill entire slide
            slide.shapes.add_picture(
                str(screenshot),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            print(f"   ✅ Added slide {i}")

            # Emit progress event
            if progress_callback:
                progress_callback('pptx_slide_added', {
                    'slide_number': i,
                    'total_slides': len(screenshots)
                })

        except Exception as e:
            print(f"   ❌ Error adding slide {i}: {str(e)}")

    # Save presentation
    prs.save(str(output_path))

    print(f"\n✅ PPTX created")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   File size: {output_path.stat().st_size / 1024:.1f} KB\n")

    return str(output_path)
